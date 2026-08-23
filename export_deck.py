"""Native PowerPoint export — a client-ready deck built from live analysis.

EVERY SLIDE IS NATIVE, EDITABLE POWERPOINT. Text is in real text frames
and numbers are in real table cells, not flattened into a picture of a
report. An advisor can retype a sentence, restyle a heading or delete a
row before presenting. That is the whole point of generating a deck
rather than attaching a PDF: a PDF is something you send, a deck is
something you finish.

Charts are the one exception, necessarily — a Plotly figure has no
PowerPoint-native equivalent, so it is rendered to a high-resolution PNG
and embedded as a picture. Rendered at 2x scale so it stays sharp on a
projector rather than pixelating.

NOTHING LEAVES THE MACHINE. No conversion service, no upload, no
third-party OCR. python-pptx writes the file locally and kaleido renders
the images locally.

THE NARRATIVE IS THE EXISTING ONE, NOT A NEW ONE. Slide 2 is built from
executive_digest.collect_flags — the same ranked strengths and concerns
the Executive Digest already shows on screen. A deck that told a
different story from the app would be worse than no deck, and a second
narrative engine is a second thing to keep true.

NUMBERS ARE NEVER INVENTED. A metric that could not be computed renders
as "not reported", never as 0.00 or a blank cell that reads as zero —
the same rule the rest of this app follows. A deck is the most likely
artefact to be shown to someone who cannot check it, so this matters
more here than anywhere else, not less.

THE DISCLOSURE TRAVELS WITH THE DECK. Every export carries the
not-investment-advice notice on the title slide. It is not optional and
not brandable: the deck leaves the building, and it must not arrive
somewhere looking like advice.
"""
import datetime
import io
import logging
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Sequence, Tuple

from logging_setup import get_logger, log_event, log_exception

logger = get_logger("export_deck")

# 16:9. python-pptx defaults to 4:3, which looks like a deck from 2003.
_SLIDE_WIDTH_IN = 13.333
_SLIDE_HEIGHT_IN = 7.5

_BLANK_LAYOUT = 6

# Rendered at 2x so charts stay sharp projected. Beyond 2x the file grows
# faster than the visible quality does.
_CHART_SCALE = 2
_CHART_PX = (1200, 620)

DISCLOSURE = (
    "Not investment advice. Generated from public market data for research purposes. "
    "Figures shown as “not reported” were unavailable, never zero."
)


@dataclass(frozen=True)
class Metric:
    """One row of the health table. `value` of None means genuinely
    unavailable, and renders as such rather than as a zero."""
    label: str
    value: Optional[float] = None
    suffix: str = ""
    decimals: int = 2

    @property
    def display(self) -> str:
        if self.value is None:
            return "not reported"
        return f"{self.value:,.{self.decimals}f}{self.suffix}"


@dataclass(frozen=True)
class DeckData:
    """Everything a deck needs, already computed by the app.

    Deliberately a plain data object with no analysis of its own: this
    module renders, it does not decide. Every value here comes from the
    same functions that produced what is on screen.
    """
    ticker: str
    company_name: str = ""
    sector: str = ""
    as_of: Optional[datetime.date] = None

    alignment_verdict: str = ""
    alignment_pct: Optional[float] = None
    strengths: Tuple[str, ...] = ()
    concerns: Tuple[str, ...] = ()

    current_price: Optional[float] = None
    intrinsic_price: Optional[float] = None
    margin_of_safety_pct: Optional[float] = None
    dcf_status: str = ""
    dcf_unavailable_reason: str = ""
    wacc: Optional[float] = None

    altman_z: Optional[float] = None
    altman_verdict: str = ""
    risk_grade: str = ""
    metrics: Tuple[Metric, ...] = ()

    # (slide title, PNG bytes)
    charts: Tuple[Tuple[str, bytes], ...] = ()

    @property
    def headline(self) -> str:
        name = self.company_name or self.ticker
        return f"{self.ticker} — {name}" if self.company_name else self.ticker

    @property
    def dcf_available(self) -> bool:
        return self.intrinsic_price is not None


def chart_png(fig, width: int = None, height: int = None,
              scale: int = _CHART_SCALE) -> Optional[bytes]:
    """A Plotly figure as PNG bytes, or None if rendering isn't possible.

    Returns None rather than raising: kaleido needs a working browser
    binary, and a machine without one should lose the chart slides, not
    the whole export. The caller reports which charts were dropped.
    """
    width = width or _CHART_PX[0]
    height = height or _CHART_PX[1]

    # Pin the figure to the export look before rendering. The app styles
    # its charts from whatever theme the user is currently viewing in, and
    # an export is not a screenshot of that: with the app in light mode the
    # axis labels come out dark and would be invisible on these slides.
    # Copied first — mutating the caller's figure would restyle the chart
    # still on screen behind the export button.
    try:
        import copy

        from export_theme import plotly_overrides

        fig = copy.deepcopy(fig)
        fig.update_layout(**plotly_overrides())
    except Exception:
        # A figure that cannot be restyled is still worth exporting.
        pass

    try:
        return fig.to_image(format="png", width=width, height=height, scale=scale)
    except Exception:
        # Kaleido serialises the figure with orjson, which handles only
        # plain JSON types — a pandas Timestamp on a date axis raises
        # "Type is not JSON serializable: Timestamp". Every price chart in
        # this app has exactly that, so the main chart failed to export
        # while the numeric-only gauges succeeded, which is a confusing way
        # to lose a slide. Plotly's OWN encoder does know these types, so a
        # round-trip through it rewrites them as ISO strings. Tried second
        # rather than always, because the round-trip costs real time on a
        # chart with a year of bars and most figures never need it.
        try:
            import json

            import plotly.graph_objects as go
            import plotly.io as pio

            normalised = go.Figure(json.loads(pio.to_json(fig)))
            return normalised.to_image(format="png", width=width, height=height,
                                       scale=scale)
        except Exception:
            log_exception(logger, "export_deck.chart_render_failed", section="export_deck")
            return None


def is_available() -> Tuple[bool, Optional[str]]:
    """Whether a deck can be built here. Returns (ok, reason).

    python-pptx is pure Python and effectively always importable; kaleido
    ships a browser binary and is the part that realistically fails, so
    they are reported separately — "install kaleido" is actionable in a
    way that "export unavailable" is not.
    """
    try:
        import pptx  # noqa: F401
    except Exception:
        return False, (
            "PowerPoint export needs the python-pptx package. Run "
            "`pip install -r requirements.txt` and restart Streamlit."
        )
    return True, None


def charts_available() -> bool:
    """Whether chart images can be rendered. A deck without charts is
    still a useful deck, so this is checked separately."""
    try:
        import plotly.graph_objects as go
        go.Figure().to_image(format="png", width=10, height=10)
        return True
    except Exception:
        return False


# --- slide construction -------------------------------------------------------

def _theme():
    """Brand name plus the shared export palette (see export_theme.py)."""
    from branding import brand
    from export_theme import palette

    return brand().name, palette()


def _paint(slide, colours):
    """Fill the slide background.

    A blank PowerPoint layout is white, and white is not the absence of a
    background — it is a background that fights every light-coloured run
    on top of it. Set explicitly on each slide rather than on the master,
    because the deck is built from the blank layout and a master edit
    would not survive a recipient changing the theme.
    """
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colours.background)
    return slide


def _blank(prs, colours):
    slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
    return _paint(slide, colours)


def _add_text(slide, text, left, top, width, height, size=18, bold=False,
              colour=None, align=None, wrap=True):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = wrap
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    # Default to the palette's body colour: PowerPoint's own default is
    # black, which is invisible on these slides.
    from export_theme import palette

    run.font.color.rgb = RGBColor.from_string(colour or palette().text)
    if align is not None:
        paragraph.alignment = align
    return box


def _add_bullets(slide, items, left, top, width, height, size=16, colour=None,
                 bullet_char="•"):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = f"{bullet_char} {item}"
        run.font.size = Pt(size)
        from export_theme import palette

        run.font.color.rgb = RGBColor.from_string(colour or palette().text)
        paragraph.space_after = Pt(8)
    return box


def _title_slide(prs, data: DeckData, brand_name: str, colours):
    from pptx.util import Inches

    slide = _blank(prs, colours)

    # The logo, when we have one. The DARK-ground file, because these
    # slides are painted near-black — the white-ground variant would sit
    # on the title slide as a white rectangle. Falls back to the wordmark
    # in text if the asset is missing, so the slide is never headless.
    _logo_placed = False
    try:
        import brand_assets

        logo = brand_assets.dark_logo()
        if logo is not None:
            # Square artwork: give it a height and let width follow, or
            # python-pptx stretches it to whatever box it is handed.
            # 1.7in, not 1.15: only about a fifth of that height is the
            # wordmark, so the smaller size left "QUANTIX" unreadable on a
            # projected slide — the same problem the tear sheet had.
            slide.shapes.add_picture(str(logo), Inches(0.8), Inches(0.45),
                                     height=Inches(1.7))
            _logo_placed = True
    except Exception:
        log_exception(logger, "export_deck.logo_failed", section="export_deck")

    if not _logo_placed:
        _add_text(slide, brand_name, 0.8, 0.7, 6, 0.5, size=16, bold=True,
                  colour=colours.accent)

    _add_text(slide, data.headline, 0.8, 2.35, 11.5, 1.2, size=40, bold=True)

    subtitle_bits = [b for b in (data.sector, data.alignment_verdict.title()) if b]
    if data.alignment_pct is not None:
        subtitle_bits.append(f"{data.alignment_pct:.0f}% alignment")
    if subtitle_bits:
        _add_text(slide, "  ·  ".join(subtitle_bits), 0.8, 3.55, 11.5, 0.5, size=18)

    as_of = data.as_of or datetime.date.today()
    _add_text(slide, f"Prepared {as_of:%d %B %Y} · Generated by {brand_name}",
              0.8, 4.15, 8, 0.4, size=14, colour=colours.text_muted)

    # The disclosure ships on the title slide, where it cannot be missed
    # by someone flicking to the numbers.
    _add_text(slide, DISCLOSURE, 0.8, 6.4, 11.5, 0.8, size=11,
              colour=colours.text_muted)
    return slide


def _synthesis_slide(prs, data: DeckData, colours):
    """Strengths and concerns straight from the Executive Digest, so the
    deck says exactly what the app says."""
    slide = _blank(prs, colours)
    _add_text(slide, "Executive Synthesis", 0.8, 0.5, 11.5, 0.7, size=30, bold=True, colour=colours.accent)

    if data.alignment_verdict:
        headline = f"{data.alignment_verdict.title()} alignment"
        if data.alignment_pct is not None:
            headline += f" — {data.alignment_pct:.0f}% of evaluable checks passed"
        _add_text(slide, headline, 0.8, 1.3, 11.5, 0.5, size=18, bold=True)

    _add_text(slide, "Strengths", 0.8, 2.1, 5.5, 0.4, size=18, bold=True, colour=colours.positive)
    _add_bullets(slide, data.strengths or ("No standout strengths flagged.",),
                 0.8, 2.6, 5.5, 3.4)

    _add_text(slide, "Concerns", 7.0, 2.1, 5.5, 0.4, size=18, bold=True, colour=colours.negative)
    _add_bullets(slide, data.concerns or ("No material concerns flagged.",),
                 7.0, 2.6, 5.5, 3.4)

    _add_text(slide,
              "Strengths and concerns are the same ranked signals shown in the app's Executive "
              "Digest. A neutral reading contributes nothing to either column.",
              0.8, 6.4, 11.5, 0.6, size=11, colour=colours.text_muted)
    return slide


def _valuation_slide(prs, data: DeckData, colours):
    slide = _blank(prs, colours)
    _add_text(slide, "Valuation", 0.8, 0.5, 11.5, 0.7, size=30, bold=True, colour=colours.accent)

    if not data.dcf_available:
        _add_text(slide,
                  data.dcf_unavailable_reason
                  or "A discounted cash flow valuation could not be computed for this company.",
                  0.8, 1.5, 11.5, 1.0, size=16)
        _add_text(slide,
                  "No intrinsic value is shown rather than an estimate built on missing inputs.",
                  0.8, 6.4, 11.5, 0.6, size=11, colour=colours.text_muted)
        return slide

    cards = [
        ("Market price", Metric("", data.current_price, decimals=2).display),
        ("Intrinsic value", Metric("", data.intrinsic_price, decimals=2).display),
        ("Margin of safety", Metric("", data.margin_of_safety_pct, suffix="%", decimals=1).display),
    ]
    for index, (label, value) in enumerate(cards):
        left = 0.8 + index * 4.1
        _add_text(slide, label, left, 1.6, 3.8, 0.4, size=14)
        _add_text(slide, value, left, 2.1, 3.8, 0.9, size=32, bold=True, colour=colours.accent)

    if data.dcf_status:
        _add_text(slide, data.dcf_status, 0.8, 3.4, 11.5, 0.6, size=18, bold=True)
    if data.wacc is not None:
        _add_text(slide, f"Discounted at a {data.wacc * 100:.2f}% weighted average cost of capital.",
                  0.8, 4.1, 11.5, 0.5, size=14)

    _add_text(slide,
              "A DCF is a model, not a measurement: it is only as good as its growth and "
              "discount-rate assumptions, which are shown in the app.",
              0.8, 6.4, 11.5, 0.6, size=11, colour=colours.text_muted)
    return slide


def _health_slide(prs, data: DeckData, colours):
    from pptx.util import Inches, Pt

    slide = _blank(prs, colours)
    _add_text(slide, "Financial Health & Risk", 0.8, 0.5, 11.5, 0.7,
              size=30, bold=True, colour=colours.accent)

    header_bits = []
    if data.altman_z is not None:
        header_bits.append(f"Altman Z {data.altman_z:.2f}"
                           + (f" ({data.altman_verdict})" if data.altman_verdict else ""))
    if data.risk_grade:
        header_bits.append(f"Risk grade {data.risk_grade}")
    if header_bits:
        _add_text(slide, "  ·  ".join(header_bits), 0.8, 1.3, 11.5, 0.5, size=18, bold=True)

    rows = list(data.metrics)
    if rows:
        # A real table, so an advisor can edit or delete a row.
        table_height = min(4.2, 0.4 + 0.36 * (len(rows) + 1))
        shape = slide.shapes.add_table(
            len(rows) + 1, 2, Inches(0.8), Inches(2.0), Inches(6.5), Inches(table_height))
        table = shape.table
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        for index, metric in enumerate(rows, start=1):
            table.cell(index, 0).text = metric.label
            table.cell(index, 1).text = metric.display

        # PowerPoint's default table style is a blue-and-white banded
        # theme. Left alone it puts white cells on a black slide — the one
        # element that would still look like it came from a different
        # document. Every cell is filled explicitly, which overrides the
        # style without having to swap the style GUID.
        from pptx.dml.color import RGBColor

        for row_index, row in enumerate(table.rows):
            header = row_index == 0
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(
                    colours.surface_alt if header else colours.surface)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
                        run.font.bold = header
                        run.font.color.rgb = RGBColor.from_string(
                            colours.accent if header else colours.text)

    _add_text(slide,
              "Values shown as “not reported” were unavailable for this company — "
              "commonly the case for financials, where several corporate ratios do not apply. "
              "They are not zeros.",
              0.8, 6.4, 11.5, 0.6, size=11, colour=colours.text_muted)
    return slide


def _chart_slide(prs, title: str, png: bytes, colours):
    from pptx.util import Inches

    slide = _blank(prs, colours)
    _add_text(slide, title, 0.8, 0.5, 11.5, 0.7, size=28, bold=True, colour=colours.accent)
    slide.shapes.add_picture(io.BytesIO(png), Inches(0.8), Inches(1.4), width=Inches(11.7))
    return slide


# --- assembly -----------------------------------------------------------------

def build_deck(data: DeckData) -> Tuple[Optional[bytes], Optional[str]]:
    """The finished .pptx as bytes. Returns (data, error); one is None.

    Never raises — an export button that throws a traceback into the page
    is worse than one that explains why it could not run.
    """
    ok, reason = is_available()
    if not ok:
        return None, reason

    try:
        from pptx import Presentation
        from pptx.util import Inches

        brand_name, colours = _theme()

        prs = Presentation()
        prs.slide_width = Inches(_SLIDE_WIDTH_IN)
        prs.slide_height = Inches(_SLIDE_HEIGHT_IN)

        _title_slide(prs, data, brand_name, colours)
        _synthesis_slide(prs, data, colours)
        _valuation_slide(prs, data, colours)
        _health_slide(prs, data, colours)
        for title, png in data.charts:
            if png:
                _chart_slide(prs, title, png, colours)

        buffer = io.BytesIO()
        prs.save(buffer)
        payload = buffer.getvalue()
        log_event(logger, logging.INFO, "export_deck.built",
                  ticker=data.ticker, slides=len(prs.slides),
                  charts=len(data.charts), bytes=len(payload))
        return payload, None
    except Exception as e:
        log_exception(logger, "export_deck.build_failed", section="export_deck")
        return None, f"Couldn't build the deck ({type(e).__name__})."


def filename_for(ticker: str, as_of: Optional[datetime.date] = None) -> str:
    as_of = as_of or datetime.date.today()
    from branding import brand
    slug = brand().name.lower().replace(" ", "-")
    return f"{slug}-{ticker.lower()}-{as_of:%Y%m%d}.pptx"
